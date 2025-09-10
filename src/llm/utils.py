import json
import re
import io
from typing import List, Dict
from pptx import Presentation



def extract_first_json(text: str) -> dict:
    """Возвращает первый JSON-объект {...}. Чинит висячие запятые при необходимости."""
    
    m = re.search(r"\{.*\}", text, flags=re.DOTALL)
    if not m:
        raise ValueError("Модель не вернула JSON-объект.")
    snippet = m.group(0)
    try:
        return json.loads(snippet)
    except json.JSONDecodeError:
        fixed = re.sub(r",\s*}", "}", re.sub(r",\s*]", "]", snippet))
        return json.loads(fixed)


def tokenize_ru(text: str) -> List[str]:
    return re.findall(r"[А-Яа-яA-Za-zЁё0-9\-]+", text.lower())


def split_sentences(text: str) -> List[str]:
    return [s for s in re.split(r'(?<=[\.\!\?])\s+', text.strip()) if s]


def get_slides_data(pptx_bytes: io.BytesIO) -> List[Dict[str, int | str]]:

    prs = Presentation(pptx_bytes)
    slides_data = []
    for i, slide in enumerate(prs.slides, start=1):
        text = []
        img_count = 0
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text.strip())
            if shape.shape_type == 13:  # PICTURE
                img_count += 1
        slides_data.append({
            "slide": i,
            "text": " ".join(text),
            "word_count": sum(len(t.split()) for t in text),
            "img_count": img_count
        })
    return slides_data